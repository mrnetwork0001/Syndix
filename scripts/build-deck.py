#!/usr/bin/env python3
"""
Builds the Syndix pitch deck as a .pptx.

Written as a script rather than a hand-made file so the deck can be regenerated
when a number changes: live chain figures come from /tmp/syndix_live.json, and
every claim in here is one that appears in the README or on /protocol. A deck
that drifts from the honesty table is worse than no deck.

Design follows the product's own tokens, so the slides and the app look like the
same thing. Run:  python3 scripts/build-deck.py
"""

import json
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- tokens

VOID = RGBColor(0x0B, 0x0B, 0x0C)
SURFACE = RGBColor(0x14, 0x14, 0x16)
ELEVATED = RGBColor(0x1A, 0x1A, 0x1E)
ACCENT = RGBColor(0x00, 0x66, 0xFF)
ACCENT_LT = RGBColor(0x4D, 0x92, 0xFF)
INK = RGBColor(0xF4, 0xF4, 0xF5)
INK_MUTED = RGBColor(0xA1, 0xA1, 0xAA)
INK_FAINT = RGBColor(0x6B, 0x6B, 0x74)
POSITIVE = RGBColor(0x22, 0xC5, 0x5E)
CAUTION = RGBColor(0xF5, 0x9E, 0x0B)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
HAIRLINE = RGBColor(0x2A, 0x2A, 0x30)

SANS = "Inter"          # Geist is not a system font; Inter is the closest match
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)   # 16:9
M = Inches(0.86)                     # outer margin

LIVE = json.load(open("/tmp/syndix_live.json")) if os.path.exists(
    "/tmp/syndix_live.json") else {
    "articleCount": 10, "uniqueReaders": 3,
    "balanceEth": 0.00595, "reservedEth": 0.00237, "block": 32093340}


# ---------------------------------------------------------------- helpers

def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = VOID
    return s


def box(s, x, y, w, h, text, size=18, color=INK, bold=False, font=SANS,
        align=PP_ALIGN.LEFT, spacing=1.25, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def rect(s, x, y, w, h, fill=SURFACE, line=HAIRLINE, line_w=Pt(0.75)):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.06
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = line_w
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def eyebrow(s, text, y=None):
    box(s, M, y or Inches(0.62), Inches(9), Inches(0.3),
        text.upper(), size=11, color=INK_FAINT, font=SANS)


def title(s, text, y=Inches(1.05), size=40, color=INK):
    box(s, M, y, W - 2 * M, Inches(1.5), text, size=size, bold=True, color=color)


def footer(s, n):
    box(s, M, H - Inches(0.62), Inches(6), Inches(0.3),
        "Syndix · GIWA GASOK", size=10, color=INK_FAINT, font=MONO)
    box(s, W - M - Inches(1.2), H - Inches(0.62), Inches(1.2), Inches(0.3),
        f"{n:02d}", size=10, color=INK_FAINT, font=MONO, align=PP_ALIGN.RIGHT)


def stat(s, x, y, w, value, label, color=INK, vsize=34):
    h = Inches(1.5)
    rect(s, x, y, w, h, fill=SURFACE)
    box(s, x + Inches(0.28), y + Inches(0.26), w - Inches(0.5), Inches(0.6),
        value, size=vsize, bold=True, color=color, font=MONO)
    box(s, x + Inches(0.28), y + Inches(0.98), w - Inches(0.5), Inches(0.4),
        label.upper(), size=10, color=INK_FAINT)


def bullets(s, x, y, w, items, size=16, gap=Inches(0.82), lead_color=INK):
    """items: list of (bold lead, rest)."""
    for i, (lead, rest) in enumerate(items):
        yy = y + i * gap
        tb = s.shapes.add_textbox(x, yy, w, gap)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.line_spacing = 1.35
        r1 = p.add_run()
        r1.text = lead + "  "
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = lead_color
        r1.font.name = SANS
        r2 = p.add_run()
        r2.text = rest
        r2.font.size = Pt(size)
        r2.font.color.rgb = INK_MUTED
        r2.font.name = SANS


def accent_bar(s, y=Inches(1.02), w=Inches(0.72)):
    sh = rect(s, M, y, w, Inches(0.055), fill=ACCENT, line=None)
    return sh


# ---------------------------------------------------------------- slides

prs = new_deck()
n = 0


def nxt():
    global n
    n += 1
    return n


# 01 — title
s = slide(prs)
rect(s, Inches(0), Inches(0), W, Inches(0.09), fill=ACCENT, line=None)
box(s, M, Inches(2.35), Inches(11), Inches(0.4), "SYNDIX",
    size=13, color=INK_FAINT, font=MONO)
box(s, M, Inches(2.85), Inches(11.2), Inches(2.0),
    "The newsroom that pays\nits readers.", size=52, bold=True, color=INK,
    spacing=1.05)
box(s, M, Inches(4.95), Inches(9.2), Inches(1.0),
    "An autonomous AI news syndicate on GIWA L2. Agents write each issue from live\n"
    "chain state. Every verified human who reads one is paid in ETH.",
    size=17, color=INK_MUTED, spacing=1.4)
box(s, M, Inches(6.25), Inches(11), Inches(0.4),
    f"GIWA GASOK  ·  Track 03 GIWA-Native Ideas  ·  chain 91342  ·  syndix.xyz",
    size=12, color=INK_FAINT, font=MONO)

# 02 — problem
s = slide(prs)
eyebrow(s, "The problem")
accent_bar(s)
title(s, "Attention is the product.\nThe reader never sees a cent.", y=Inches(1.35), size=38)
box(s, M, Inches(3.5), Inches(10.6), Inches(1.4),
    "Publishing monetises attention by selling it to advertisers. The reader produces "
    "the value and receives none of it.\n\n"
    "Paying readers directly has been impractical for two reasons, and both are "
    "properties of the chain rather than the product.",
    size=17, color=INK_MUTED, spacing=1.45)
rect(s, M, Inches(5.35), Inches(5.4), Inches(1.15), fill=SURFACE)
box(s, M + Inches(0.32), Inches(5.6), Inches(4.8), Inches(0.8),
    "1.  A meaningful micro-reward costs\n     more to send than it is worth.",
    size=14, color=INK, spacing=1.35)
rect(s, M + Inches(5.75), Inches(5.35), Inches(5.4), Inches(1.15), fill=SURFACE)
box(s, M + Inches(6.07), Inches(5.6), Inches(4.8), Inches(0.8),
    "2.  An open reward pool is drained by\n     scripts faster than humans can read.",
    size=14, color=INK, spacing=1.35)
footer(s, nxt() + 1)

# 03 — why GIWA
s = slide(prs)
eyebrow(s, "Why this is only viable on GIWA")
accent_bar(s)
title(s, "The reward is 166× the gas\nneeded to deliver it.", y=Inches(1.35), size=38)
y = Inches(3.6)
stat(s, M, y, Inches(3.55), "0.00003", "ETH reward per reader", color=INK)
stat(s, M + Inches(3.85), y, Inches(3.55), "0.00000018", "ETH gas to deliver it", color=INK)
stat(s, M + Inches(7.70), y, Inches(3.55), "166×", "reward-to-cost ratio", color=ACCENT_LT)
box(s, M, Inches(5.5), Inches(11.2), Inches(1.4),
    "On Ethereum L1 that ratio inverts and the product cannot exist. GIWA's ~1s blocks and "
    "sub-cent fees are not a convenience here, they are the precondition.\n\n"
    "Flashblocks confirm the claim in about 200ms, so the reward lands while the reader is "
    "still looking at the button.",
    size=15, color=INK_MUTED, spacing=1.4)
footer(s, nxt() + 1)

# 04 — sybil
s = slide(prs)
eyebrow(s, "The second blocker: sybil resistance")
accent_bar(s)
title(s, "We cannot issue the\ncredential we check.", y=Inches(1.35), size=38)
box(s, M, Inches(2.75), Inches(11.2), Inches(2.2),
    "The gate is Upbit Web3 Names — soul-bound, one per verified wallet, issued by GIWA "
    "through Dojang attestation. Syndix reads it and cannot write it.\n\n"
    "We hold syndix.up.id, obtained the same way any reader does. What we cannot do is "
    "mint one - for a reader, or for ourselves. A protocol that can issue itself the "
    "credential it gates on has not built a gate.",
    size=17, color=INK_MUTED, spacing=1.45)
rect(s, M, Inches(5.15), Inches(11.2), Inches(1.25), fill=ELEVATED, line=ACCENT)
box(s, M + Inches(0.4), Inches(5.42), Inches(10.4), Inches(0.8),
    "UpIdReaderRegistry.isVerified(wallet)  →  UPNAME.balanceOf(wallet) > 0\n"
    "One human, one name, one claim per issue.",
    size=14, color=INK, font=MONO, spacing=1.5)
footer(s, nxt() + 1)

# 05 — how it works
s = slide(prs)
eyebrow(s, "How it works")
accent_bar(s)
title(s, "A newsroom, not a platform.", y=Inches(1.35), size=38)
box(s, M, Inches(2.5), Inches(11), Inches(0.5),
    "Nobody submits an article, so nothing is moderated.",
    size=16, color=INK_MUTED)
steps = [
    ("SCAN", "Read GIWA head state and gas price from the Flashblocks RPC."),
    ("GENERATE", "gpt-4.1 writes the issue against a strict JSON schema."),
    ("PIN", "Body goes to IPFS. Publishing is blocked if pinning fails."),
    ("PUBLISH", "publishArticle records it; the attached ETH is the reward pool."),
    ("READ & CLAIM", "The reader proves attention, then submits their own claim."),
]
yy = Inches(3.25)
for i, (k, v) in enumerate(steps):
    y = yy + i * Inches(0.76)
    rect(s, M, y, Inches(11.2), Inches(0.62), fill=SURFACE)
    box(s, M + Inches(0.3), y + Inches(0.17), Inches(1.9), Inches(0.35),
        f"{i+1}. {k}", size=12, bold=True, color=ACCENT_LT, font=MONO)
    box(s, M + Inches(2.5), y + Inches(0.17), Inches(8.4), Inches(0.35),
        v, size=13, color=INK_MUTED)
footer(s, nxt() + 1)

# 06 — proof of read
s = slide(prs)
eyebrow(s, "The hard problem")
accent_bar(s)
title(s, "Proving the attention\nactually happened.", y=Inches(1.35), size=38)
box(s, M, Inches(3.35), Inches(5.2), Inches(2.4),
    "Client-reported dwell time is worthless. A claim would cost one HTTP request "
    "with no page load at all.\n\n"
    "So the server keeps the clock. A signed session is stamped on open, the page "
    "beats with its scroll depth, and beats arriving faster than real time are refused.",
    size=15, color=INK_MUTED, spacing=1.45)
rows = [
    ("No session", "refused"),
    ("Edited token", "refused, HMAC fails"),
    ("Spammed heartbeats", "too soon"),
    ("Full scroll, no time", "refused"),
    ("Full time, no scroll", "refused"),
    ("Session from another issue", "refused"),
]
x2 = M + Inches(5.75)
rect(s, x2, Inches(3.3), Inches(5.45), Inches(3.1), fill=SURFACE)
for i, (a, b) in enumerate(rows):
    y = Inches(3.52) + i * Inches(0.46)
    box(s, x2 + Inches(0.3), y, Inches(3.1), Inches(0.35), a, size=12, color=INK_MUTED)
    box(s, x2 + Inches(3.4), y, Inches(1.9), Inches(0.35), b, size=12,
        color=POSITIVE, font=MONO, align=PP_ALIGN.RIGHT)
box(s, M, Inches(6.05), Inches(11.2), Inches(0.6),
    "Stated plainly: this proves time and scrolling, not comprehension. It prevents instant "
    "and bulk claiming, and up.id caps a human at one claim per issue.",
    size=13, color=INK_FAINT, spacing=1.4)
footer(s, nxt() + 1)

# 07 — what is live
s = slide(prs)
eyebrow(s, "Shipped, on testnet, today")
accent_bar(s)
title(s, "Not a prototype.", y=Inches(1.35), size=38)
y = Inches(2.85)
stat(s, M, y, Inches(2.6), str(LIVE["articleCount"]), "issues published onchain")
stat(s, M + Inches(2.9), y, Inches(2.6), str(LIVE["uniqueReaders"]), "wallets paid")
stat(s, M + Inches(5.8), y, Inches(2.6), "101", "Foundry tests", color=ACCENT_LT)
stat(s, M + Inches(8.7), y, Inches(2.5), "5", "verified contracts", color=ACCENT_LT)
contracts = [
    ("SyndixTreasury", "pools, claims, solvency"),
    ("SyndixArticleNFT", "two-level open edition"),
    ("UpIdReaderRegistry", "sybil gate over up.id"),
    ("SyndixPaymaster", "ERC-4337 v0.7, staked and funded"),
]
yy = Inches(4.7)
for i, (a, b) in enumerate(contracts):
    y = yy + i * Inches(0.44)
    box(s, M, y, Inches(3.6), Inches(0.35), a, size=13, color=INK, font=MONO)
    box(s, M + Inches(3.8), y, Inches(7), Inches(0.35), b, size=13, color=INK_FAINT)
box(s, M, Inches(6.55), Inches(11.2), Inches(0.4),
    f"All verified on the GIWA explorer  ·  read at block {LIVE['block']:,}",
    size=12, color=INK_FAINT, font=MONO)
footer(s, nxt() + 1)

# 08 — business model
s = slide(prs)
eyebrow(s, "Sustainability")
accent_bar(s)
title(s, "We do not invent money\nto pay readers.", y=Inches(1.35), size=38)
box(s, M, Inches(3.4), Inches(11.2), Inches(1.2),
    "We redirect the advertising budget one step further down the chain. Today an advertiser "
    "pays a publisher to harvest a reader's attention and the reader gets nothing. Here the "
    "sponsor's money lands in the reader's wallet and the protocol takes a capped fee.",
    size=17, color=INK_MUTED, spacing=1.45)
rect(s, M, Inches(4.85), Inches(5.4), Inches(1.55), fill=SURFACE)
box(s, M + Inches(0.32), Inches(5.08), Inches(4.8), Inches(1.1),
    "It is a better ad product\nProof of read is onchain, so there are no phantom "
    "impressions. One up.id is one human, so there is no bot traffic.",
    size=13, color=INK_MUTED, spacing=1.35)
rect(s, M + Inches(5.8), Inches(4.85), Inches(5.4), Inches(1.55), fill=SURFACE)
box(s, M + Inches(6.12), Inches(5.08), Inches(4.8), Inches(1.1),
    "The sponsor can verify it\nSyndixSponsorship splits each deposit into a capped fee "
    "and a reader share no owner function can reach.",
    size=13, color=INK_MUTED, spacing=1.35)
footer(s, nxt() + 1)

# 09 — unit economics
s = slide(prs)
eyebrow(s, "Unit economics")
accent_bar(s)
title(s, "The margin arrives with the audience.", y=Inches(1.35), size=34)
hdr = ["", "Today, 20 readers", "At 1,000 readers"]
rowdata = [
    ("Reward cost per issue", "$1.14", "$57"),
    ("Generation + gas", "~$0.03", "~$3"),
    ("Sponsor price", "n/a", "$150 – $300"),
    ("Margin per issue", "subsidised", "$90 – $240"),
]
y0 = Inches(3.0)
rect(s, M, y0, Inches(11.2), Inches(0.55), fill=ELEVATED)
for j, htxt in enumerate(hdr):
    box(s, M + Inches(0.3) + j * Inches(3.7), y0 + Inches(0.16),
        Inches(3.5), Inches(0.35), htxt.upper(), size=10, color=INK_FAINT)
for i, row in enumerate(rowdata):
    y = y0 + Inches(0.55) + i * Inches(0.62)
    rect(s, M, y, Inches(11.2), Inches(0.62), fill=SURFACE)
    for j, cell in enumerate(row):
        col = INK if j == 0 else (POSITIVE if i == 3 and j == 2 else INK_MUTED)
        fnt = SANS if j == 0 else MONO
        box(s, M + Inches(0.3) + j * Inches(3.7), y + Inches(0.19),
            Inches(3.5), Inches(0.35), cell, size=13, color=col, font=fnt)
box(s, M, Inches(6.15), Inches(11.2), Inches(0.7),
    "Newsletter sponsorships run $25–50 CPM against unverified audiences. Proven attention "
    "justifies the top of that band. The grant is the runway to build the audience that makes "
    "it work.",
    size=13, color=INK_FAINT, spacing=1.4)
footer(s, nxt() + 1)

# 10 — honesty
s = slide(prs)
eyebrow(s, "What is real and what is not")
accent_bar(s)
title(s, "We publish our own gaps.", y=Inches(1.35), size=38)
box(s, M, Inches(2.5), Inches(11), Inches(0.4),
    "Every surface in the app that shows something other than live chain data says so. "
    "The same table is on syndix.xyz/protocol.",
    size=14, color=INK_MUTED)
real = ["Smart contracts", "Reader reward claim", "Proof of read", "up.id identity",
        "Issue content + IPFS", "Analytics from event logs", "x402 settlement"]
notyet = ["Autonomous publishing", "Sponsorship revenue", "Gasless via ERC-4337",
          "KRW denomination"]
rect(s, M, Inches(3.2), Inches(5.4), Inches(3.2), fill=SURFACE)
box(s, M + Inches(0.3), Inches(3.42), Inches(4.8), Inches(0.3),
    "REAL TODAY", size=11, bold=True, color=POSITIVE)
for i, t in enumerate(real):
    box(s, M + Inches(0.3), Inches(3.82) + i * Inches(0.33), Inches(4.8),
        Inches(0.3), t, size=12, color=INK_MUTED)
rect(s, M + Inches(5.8), Inches(3.2), Inches(5.4), Inches(3.2), fill=SURFACE)
box(s, M + Inches(6.1), Inches(3.42), Inches(4.8), Inches(0.3),
    "NOT YET, AND SAID SO", size=11, bold=True, color=CAUTION)
for i, t in enumerate(notyet):
    box(s, M + Inches(6.1), Inches(3.82) + i * Inches(0.33), Inches(4.8),
        Inches(0.3), t, size=12, color=INK_MUTED)
box(s, M + Inches(6.1), Inches(5.35), Inches(4.8), Inches(0.9),
    "Two of these are already written and tested as contracts, awaiting deployment: "
    "SyndixPublisher and SyndixSponsorship.",
    size=11, color=INK_FAINT, spacing=1.35)
footer(s, nxt() + 1)

# 11 — roadmap
s = slide(prs)
eyebrow(s, "Roadmap")
accent_bar(s)
title(s, "Each item closes a named gap.", y=Inches(1.35), size=38)
items = [
    ("01", "Unattended newsroom",
     "Deploy SyndixPublisher and schedule the pipeline, without an owner key on a server."),
    ("02", "Sponsored issues",
     "Deploy SyndixSponsorship. Turn the treasury from a subsidy into revenue."),
    ("03", "Comprehension challenges",
     "Raise proof of read above time and scroll, using the issue's own content."),
    ("04", "Gasless claims",
     "The paymaster is deployed and funded; it waits on a bundler for chain 91342."),
    ("05", "KRW denomination",
     "SyndixStableTreasury is written and tested. A 100 KRW promise pays 100 KRW."),
]
yy = Inches(2.9)
for i, (num, head, body) in enumerate(items):
    y = yy + i * Inches(0.82)
    box(s, M, y, Inches(0.6), Inches(0.4), num, size=14, bold=True,
        color=ACCENT_LT, font=MONO)
    box(s, M + Inches(0.75), y, Inches(3.6), Inches(0.4), head, size=15,
        bold=True, color=INK)
    box(s, M + Inches(4.5), y, Inches(6.6), Inches(0.6), body, size=13,
        color=INK_MUTED, spacing=1.3)
footer(s, nxt() + 1)

# 12 — ask / close
s = slide(prs)
rect(s, Inches(0), Inches(0), W, Inches(0.09), fill=ACCENT, line=None)
eyebrow(s, "The ask", y=Inches(1.9))
box(s, M, Inches(2.35), Inches(11.2), Inches(1.6),
    "Fund the audience that makes\nthe economics work.", size=44, bold=True,
    color=INK, spacing=1.08)
box(s, M, Inches(4.25), Inches(10.4), Inches(1.2),
    "The protocol is built, deployed and honest about its gaps. What it does not have yet is "
    "readers. A GASOK grant funds roughly 350,000 reader rewards at today's rate — which is "
    "the runway from subsidy to sponsorship revenue.",
    size=16, color=INK_MUTED, spacing=1.45)
rect(s, M, Inches(5.75), Inches(11.2), Inches(0.9), fill=SURFACE)
box(s, M + Inches(0.4), Inches(5.98), Inches(10.4), Inches(0.5),
    "syndix.xyz   ·   syndix.xyz/tech   ·   github.com/mrnetwork0001/Syndix",
    size=14, color=INK, font=MONO)
footer(s, nxt() + 1)

out = "Syndix-Pitch-Deck.pptx"
prs.save(out)
print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
